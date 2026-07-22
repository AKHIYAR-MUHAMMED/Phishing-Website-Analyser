import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors matching reference template
    NAVY_DARK = RGBColor(16, 37, 66)      # Primary Dark Text / Accent
    HEADER_BLUE = RGBColor(0, 51, 102)    # Department Banner Blue
    ORANGE_RUST = RGBColor(192, 80, 0)   # Diagram Block Accent
    DARK_BROWN = RGBColor(125, 50, 0)    # Workflow Accent
    TEXT_BLACK = RGBColor(30, 30, 30)     # Main Body Text
    MUTED_GRAY = RGBColor(100, 100, 100) # Subtext / Footers
    WHITE = RGBColor(255, 255, 255)
    CARD_BG = RGBColor(235, 240, 245)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    def add_footer(slide, current_slide, total_slides=24):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(11.733), Inches(0.35))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        run1 = p.add_run()
        run1.text = "23-07-2026          "
        run1.font.size = Pt(11)
        run1.font.color.rgb = MUTED_GRAY
        run1.font.name = "Calibri"
        
        run2 = p.add_run()
        run2.text = "Department of Artificial Intelligence & Data Science, ASIET"
        run2.font.size = Pt(11)
        run2.font.color.rgb = MUTED_GRAY
        run2.font.name = "Calibri"
        
        run3 = p.add_run()
        run3.text = f"          {current_slide}"
        run3.font.size = Pt(11)
        run3.font.color.rgb = MUTED_GRAY
        run3.font.name = "Calibri"

    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(10.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text.upper()
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = TEXT_BLACK
        run.font.name = "Calibri"
        run.font.underline = True

        # Top Right Logo Badge
        logo_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.0), Inches(0.3), Inches(0.8), Inches(0.8))
        logo_shape.fill.solid()
        logo_shape.fill.fore_color.rgb = HEADER_BLUE
        logo_shape.line.color.rgb = ORANGE_RUST
        logo_shape.line.width = Pt(1.5)
        
        tf_logo = logo_shape.text_frame
        p_logo = tf_logo.paragraphs[0]
        p_logo.alignment = PP_ALIGN.CENTER
        r_logo = p_logo.add_run()
        r_logo.text = "ASIET"
        r_logo.font.size = Pt(9)
        r_logo.font.bold = True
        r_logo.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Top Left Institution Header Text Box
    inst_box = slide1.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(6.5), Inches(1.1))
    tf_inst = inst_box.text_frame
    tf_inst.margin_left = tf_inst.margin_top = tf_inst.margin_right = tf_inst.margin_bottom = 0
    p_i1 = tf_inst.paragraphs[0]
    r_i1 = p_i1.add_run()
    r_i1.text = "Adi Shankara\n"
    r_i1.font.size = Pt(22)
    r_i1.font.bold = True
    r_i1.font.color.rgb = HEADER_BLUE
    r_i1.font.name = "Georgia"
    
    r_i2 = p_i1.add_run()
    r_i2.text = "INSTITUTE OF ENGINEERING AND TECHNOLOGY\n"
    r_i2.font.size = Pt(10)
    r_i2.font.bold = True
    r_i2.font.color.rgb = TEXT_BLACK
    
    r_i3 = p_i1.add_run()
    r_i3.text = "Approved by AICTE & Affiliated to APJ Abdul Kalam Technological University\n(Owned by Adi Sankara Trust)"
    r_i3.font.size = Pt(8.5)
    r_i3.font.color.rgb = MUTED_GRAY

    # Banner Top Right Box
    dept_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(0.4), Inches(5.0), Inches(0.9))
    dept_box.fill.solid()
    dept_box.fill.fore_color.rgb = HEADER_BLUE
    dept_box.line.fill.background()
    tf_dept = dept_box.text_frame
    tf_dept.margin_left = Inches(0.2)
    p_dept = tf_dept.paragraphs[0]
    p_dept.alignment = PP_ALIGN.CENTER
    r_dept = p_dept.add_run()
    r_dept.text = "DEPARTMENT OF\nARTIFICIAL INTELLIGENCE AND DATA SCIENCE"
    r_dept.font.size = Pt(11)
    r_dept.font.bold = True
    r_dept.font.color.rgb = WHITE
    
    # Title Text Box
    title_box1 = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(1.8))
    tf1 = title_box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "PHISHING WEBSITE IMPERSONATION:\nCOMPARATIVE ANALYSIS OF DETECTION AND TARGET RECOGNITION METHODS"
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = TEXT_BLACK
    r1.font.name = "Calibri"

    # Presenter Details Box
    details_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(9.0), Inches(2.2))
    tf_det = details_box.text_frame
    
    details = [
        ("Presented By : ", "Marcin Jarczewski, Piotr Białczak, Wojciech Mazurczyk"),
        ("Reg No. : ", "ASI22CA035 / CERT Polska & WUT"),
        ("Guide : ", "Asst. Prof. Asha Rose Thomas"),
        ("Semester : ", "S7, CSE(AI)")
    ]
    for label, val in details:
        p = tf_det.add_paragraph() if tf_det.paragraphs[0].text else tf_det.paragraphs[0]
        p.space_after = Pt(4)
        r_lbl = p.add_run()
        r_lbl.text = label
        r_lbl.font.size = Pt(15)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = TEXT_BLACK
        
        r_val = p.add_run()
        r_val.text = val
        r_val.font.size = Pt(15)
        r_val.font.color.rgb = TEXT_BLACK

    add_footer(slide1, 1)

    def create_bullet_slide(slide_num, title, bullets_data, font_size_primary=16, font_size_secondary=14, space_after=6):
        slide = prs.slides.add_slide(blank_slide_layout)
        add_slide_header(slide, title)
        
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.3))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = 0
        
        for idx, item in enumerate(bullets_data):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            level = item.get("level", 0)
            p.level = level
            p.space_after = Pt(space_after)
            
            prefix = "• " if level == 0 else "o " if level == 1 else "- "
            
            r_pre = p.add_run()
            r_pre.text = prefix
            r_pre.font.size = Pt(font_size_primary if level == 0 else font_size_secondary)
            r_pre.font.bold = True
            r_pre.font.color.rgb = TEXT_BLACK
            
            if "bold_lead" in item:
                r_lead = p.add_run()
                r_lead.text = item["bold_lead"] + " "
                r_lead.font.size = Pt(font_size_primary if level == 0 else font_size_secondary)
                r_lead.font.bold = True
                r_lead.font.color.rgb = TEXT_BLACK
                
            r_text = p.add_run()
            r_text.text = item["text"]
            r_text.font.size = Pt(font_size_primary if level == 0 else font_size_secondary)
            r_text.font.color.rgb = TEXT_BLACK
            
        add_footer(slide, slide_num)
        return slide

    # =========================================================================
    # SLIDE 2: CONTENTS (Optimized font sizes for zero overlap)
    # =========================================================================
    contents_data = [
        {"bold_lead": "Cybersecurity Challenge", "text": "& Phishing Escalation"},
        {"bold_lead": "Limitations of Traditional", "text": "Heuristics & Code Analysis"},
        {"bold_lead": "The Visual Analysis Solution", "text": "& Dual Research Objectives"},
        {"bold_lead": "Key Technologies:", "text": "Deep Learning vs. Perceptual Hashing"},
        {"bold_lead": "Evaluated Method 1:", "text": "VisualPhishNet (Triplet CNN Layout)"},
        {"bold_lead": "Evaluated Method 2:", "text": "Phishpedia (Faster R-CNN + Siamese)"},
        {"bold_lead": "Proposed Baseline Method:", "text": "Perceptual Hashing (pHash + FAISS)"},
        {"bold_lead": "Experimental Framework", "text": "& Dataset Standardization"},
        {"bold_lead": "Thresholding Strategy:", "text": "Equal Error Rate (EER) Optimization"},
        {"bold_lead": "Performance Analysis:", "text": "Binary Classification & Feature Collapse"},
        {"bold_lead": "Performance Analysis:", "text": "Multiclass Target Brand Recognition"},
        {"bold_lead": "Hybrid Operational Pipeline", "text": "& Real-World Applications"},
        {"bold_lead": "Industry Impact & Future Scope", "text": ""},
        {"bold_lead": "Conclusion", "text": "& Key Takeaways"}
    ]
    create_bullet_slide(2, "CONTENTS", contents_data, font_size_primary=14, font_size_secondary=12, space_after=3)

    # Slides 3 to 16, 19 to 23
    s3_data = [
        {"bold_lead": "The Rapid Digitization of Global Commerce:", "text": "Over 72% of enterprises maintain online portals, with 27% of global sales executed digitally."},
        {"bold_lead": "Escalating Scale of Cybercrime:", "text": "Observed phishing attacks surpassed 1.13 Million in Q2 2025 (rising from 1.00M in Q1 2025)."},
        {"bold_lead": "Primary Attack Vectors & Delivery:", "text": "Fraudulent websites, email campaigns, SMS (smishing), and voice calls (vishing)."},
        {"bold_lead": "Key Problems for Security Teams:", "text": ""},
        {"bold_lead": "High Volume & Sophistication:", "text": "Attackers deploy fast-flux networks and automated site generators.", "level": 1},
        {"bold_lead": "High Economic Impact:", "text": "Credential theft, financial fraud, and severe organizational brand damage.", "level": 1},
        {"bold_lead": "Passive Defense Limits:", "text": "Static blocklists fail against zero-day phishing infrastructure.", "level": 1}
    ]
    create_bullet_slide(3, "CYBERSECURITY CHALLENGE", s3_data)

    s4_data = [
        {"bold_lead": "The Obfuscation Paradox:", "text": "Attackers heavily scramble source code while maintaining identical visual rendering for victims."},
        {"bold_lead": "Weaknesses of Legacy Detection:", "text": ""},
        {"bold_lead": "URL & Text Heuristics:", "text": "Easily bypassed via dynamic subdomains, FreeURL manipulation, and IP redirection.", "level": 1},
        {"bold_lead": "HTML Source Code Analysis:", "text": "Vulnerable to JavaScript obfuscation, DOM hiding, and dynamic content fetching.", "level": 1},
        {"bold_lead": "Early Computer Vision (EMD/SIFT/DAISY):", "text": "Computationally prohibitive for real-time traffic and lack generalization.", "level": 1},
        {"bold_lead": "The Imperative for Visual Analysis:", "text": "Regardless of code obfuscation, the visual representation must mimic the trusted brand to deceive users."}
    ]
    create_bullet_slide(4, "LIMITATIONS OF TRADITIONAL HEURISTICS", s4_data)

    s5_data = [
        {"bold_lead": "Visual Phishing Detection & Target Recognition Pipeline:", "text": "Treating webpage screenshots as the ground-truth data source for automated security triage."},
        {"bold_lead": "Dual Research Objectives:", "text": ""},
        {"bold_lead": "Binary Phishing Detection:", "text": "Determine whether a site possesses malicious intent (Phishing vs. Benign).", "level": 1},
        {"bold_lead": "Multiclass Target Recognition:", "text": "Identify the exact brand/institution being impersonated to enable specific incident response.", "level": 1},
        {"bold_lead": "The Concept in 3 Steps:", "text": ""},
        {"bold_lead": "INPUT:", "text": "Webpage screenshot collected by operational crawler / CERT analyst.", "level": 1},
        {"bold_lead": "PROCESSING:", "text": "Standardized microservices pipeline combining feature extraction & vector matching.", "level": 1},
        {"bold_lead": "OUTPUT:", "text": "Binary threat verdict + exact impersonated brand identification.", "level": 1}
    ]
    create_bullet_slide(5, "THE CORE IDEA: PROPOSED SOLUTION", s5_data)

    s6_data = [
        {"bold_lead": "The Proposed Lightweight Baseline:", "text": "Perceptual Hashing (pHash) combined with FAISS similarity search."},
        {"bold_lead": "Locality-Sensitive Hashing (LSH) Concept:", "text": "Unlike cryptographic hashes (MD5/SHA256) where minor edits alter the output, LSH maps visually similar images to close hash buckets."},
        {"bold_lead": "Discrete Cosine Transform (DCT) pHash:", "text": ""},
        {"bold_lead": "Frequency Conversion:", "text": "Transforms screenshots into frequency domains, preserving structural layout features while stripping noise.", "level": 1},
        {"bold_lead": "Vector Hashing (pHashF):", "text": "Generates real-valued vectors for dense similarity evaluation.", "level": 1},
        {"bold_lead": "FAISS Library Integration:", "text": "Enables high-speed Euclidean and Hamming distance vector indexing across millions of reference targets."}
    ]
    create_bullet_slide(6, "KEY TECHNOLOGY 1: PERCEPTUAL HASHING & FAISS", s6_data)

    s7_data = [
        {"bold_lead": "State-of-the-Art Deep Learning Approaches:", "text": "Automated neural feature extraction replacing handcrafted vision descriptors."},
        {"bold_lead": "Two Distinct Design Philosophies:", "text": ""},
        {"bold_lead": "Holistic Layout Analysis (VisualPhishNet):", "text": "Evaluates the overall visual 'feel', spatial hierarchy, and layout composition of the entire webpage.", "level": 1},
        {"bold_lead": "Targeted Object Detection (Phishpedia):", "text": "Focuses specifically on high-precision brand logo extraction and identity matching.", "level": 1},
        {"bold_lead": "The Critical Research Question:", "text": "Do complex deep neural networks consistently outperform lightweight baseline algorithms across diverse, real-world datasets?"}
    ]
    create_bullet_slide(7, "KEY TECHNOLOGY 2: DEEP LEARNING VISUAL MODELS", s7_data)

    s8_data = [
        {"bold_lead": "Holistic Layout Matching vs. Logo Object Detection:", "text": ""},
        {"bold_lead": "Holistic Approach (Layout-Based):", "text": ""},
        {"bold_lead": "Premise:", "text": "Login pages and homepages of a brand share consistent visual styling.", "level": 1},
        {"bold_lead": "Mechanism:", "text": "Convolutional Neural Networks (CNNs) encode the full screenshot into a dense embedding space.", "level": 1},
        {"bold_lead": "Targeted Approach (Logo-Based):", "text": ""},
        {"bold_lead": "Premise:", "text": "Phishing sites must display authentic brand logos to gain victim trust.", "level": 1},
        {"bold_lead": "Mechanism:", "text": "Bounding box proposal models detect logo regions and compare them against protected reference databases.", "level": 1}
    ]
    create_bullet_slide(8, "DEEP DIVE: VISUAL SIMILARITY VS OBJECT DETECTION", s8_data)

    s9_data = [
        {"bold_lead": "Phishpedia Architecture (Lin et al.):", "text": "A high-precision two-stage visual brand recognition pipeline."},
        {"bold_lead": "Stage 1: Logo Region Proposal:", "text": "Utilizes a Faster R-CNN object detection model fine-tuned on the Logo-2K+ dataset to crop potential logo bounding boxes."},
        {"bold_lead": "Stage 2: Siamese Identity Matching:", "text": "Passes cropped regions into a Siamese Network to compute distance against a protected brand reference database."},
        {"bold_lead": "Key Design Choice:", "text": "Replaces standard triplet loss with classification-based fine-tuning to effectively distinguish fine-grained brand variants (e.g., 'Adobe' vs. 'Adobe AIR')."}
    ]
    create_bullet_slide(9, "MODEL ARCHITECTURE 1: PHISHPEDIA", s9_data)

    s10_data = [
        {"bold_lead": "VisualPhishNet Architecture (Abdelnabi et al.):", "text": "A holistic webpage layout representation network."},
        {"bold_lead": "Core Backbone:", "text": "Uses a VGG-16 Convolutional Neural Network trained with Triplet Loss."},
        {"bold_lead": "Triplet Loss Mechanics:", "text": ""},
        {"bold_lead": "Inputs:", "text": "Anchor image (reference site), Positive sample (same brand page), Negative sample (different brand page).", "level": 1},
        {"bold_lead": "Objective:", "text": "Minimizes Euclidean distance between Anchor and Positive while maximizing distance to Negative by a defined margin.", "level": 1},
        {"bold_lead": "Hard Negative Mining:", "text": "Selects visually confusing non-brand pages during training to refine decision boundaries."}
    ]
    create_bullet_slide(10, "MODEL ARCHITECTURE 2: VISUALPHISHNET", s10_data)

    s11_data = [
        {"bold_lead": "Addressing Dataset Fragmentation:", "text": "Evaluated across three comprehensive datasets to eliminate single-benchmark bias."},
        {"bold_lead": "The Three Evaluated Datasets:", "text": ""},
        {"bold_lead": "CERT Polska Dataset:", "text": "15,049 real-world, analyst-verified URLs & screenshots (36 target brands, unaugmented).", "level": 1},
        {"bold_lead": "Phishpedia (PP) Dataset:", "text": "14,500 phishing / 1,542 benign screenshots (56 targets, augmented).", "level": 1},
        {"bold_lead": "VisualPhishNet (VP) Dataset:", "text": "4,644 phishing / 8,835 benign screenshots (144 targets, augmented).", "level": 1},
        {"bold_lead": "Stratified Data Splitting:", "text": "60% Training, 20% Validation, 20% Testing split applied consistently across all target classes."}
    ]
    create_bullet_slide(11, "WORKFLOW STEP 1: DATASET CURATION & STANDARDIZATION", s11_data)

    s12_data = [
        {"bold_lead": "Containerized Microservices Architecture:", "text": "Engineered to resolve 'dependency hell' across conflicting deep learning library versions."},
        {"bold_lead": "Framework Layer Breakdown:", "text": ""},
        {"bold_lead": "API Gateway (FastAPI):", "text": "Acts as central orchestrator, validating Base64 screenshot payloads via Pydantic.", "level": 1},
        {"bold_lead": "Isolated Model Containers:", "text": "Dedicated Docker containers hosting VisualPhishNet, Phishpedia, and Baseline models independently.", "level": 1},
        {"bold_lead": "Persistence Layer (SQLite & SQLAlchemy):", "text": "Captures binary predictions, target attribution labels, and confidence metrics for post-hoc analysis.", "level": 1},
        {"bold_lead": "Open Source Release:", "text": "Shared with the cybersecurity community on GitHub for reproducible benchmarking."}
    ]
    create_bullet_slide(12, "WORKFLOW STEP 2: EXPERIMENTAL MICROSERVICES FRAMEWORK", s12_data)

    s13_data = [
        {"bold_lead": "Distance-Based Decision Boundaries:", "text": "Classification relies on embedding distance to nearest protected target."},
        {"bold_lead": "Two-Stage Threshold Optimization Strategy:", "text": ""},
        {"bold_lead": "Coarse Search:", "text": "Evaluates distance thresholds from 0 to maximum validation distance in steps of 10.", "level": 1},
        {"bold_lead": "Fine-Grained Search:", "text": "Refines search within Mean ± 1 Std Dev in single-unit increments.", "level": 1},
        {"bold_lead": "Equal Error Rate (EER) Point:", "text": "Identifies the exact threshold where False Positive Rate (FPR) equals False Negative Rate (FNR)."},
        {"bold_lead": "Optimal Threshold Findings:", "text": "CERT Polska = 8.00 | Phishpedia (PP) = 3.00 | VisualPhishNet (VP) = 50.00."}
    ]
    create_bullet_slide(13, "WORKFLOW STEP 3: THRESHOLDING & EER OPTIMIZATION", s13_data)

    s14_data = [
        {"bold_lead": "Binary Detection Results (Phishing vs. Benign):", "text": ""},
        {"bold_lead": "Baseline (pHash) Superiority:", "text": "Achieved top binary performance across all datasets (PP F1 = 0.9539, VP ROC AUC = 0.8201, VP MCC = 0.6294).", "level": 1},
        {"bold_lead": "Phishpedia Binary Performance:", "text": "Strong on PP dataset (F1 = 0.9062), but struggled on CERT Polska (F1 = 0.1598) due to strict logo dependence.", "level": 1},
        {"bold_lead": "VisualPhishNet Failure & Feature Collapse:", "text": ""},
        {"bold_lead": "Low Metrics:", "text": "Achieved F1 of only 0.1673 on PP dataset and negative MCC values (-0.6160).", "level": 1},
        {"bold_lead": "Root Cause:", "text": "Overfitted to synthetic data artifacts; 85% of detected phishing samples collapsed into just 3 target clusters (Adobe, Absa, Paschoalotto).", "level": 1}
    ]
    create_bullet_slide(14, "WORKFLOW STEP 4: BINARY CLASSIFICATION & FEATURE COLLAPSE", s14_data)

    s15_data = [
        {"bold_lead": "Multiclass Impersonation Target Identification:", "text": "Evaluating exact brand attribution among 36 to 144 target classes."},
        {"bold_lead": "Phishpedia Dominance in Target Recognition:", "text": ""},
        {"bold_lead": "High Identification Rate:", "text": "Maintained >0.90 Identification Rate across ALL datasets (CERT = 0.9845, VP = 0.9270, PP = 0.9154).", "level": 1},
        {"bold_lead": "Robustness:", "text": "Logo region extraction effectively handles brand variants regardless of site layout changes.", "level": 1},
        {"bold_lead": "Baseline Multiclass Performance:", "text": "Achieved solid F1 Micro (VP = 0.7009) but lower Identification Rate (0.2554 - 0.5687), proving better suited for binary filtering than exact brand attribution."}
    ]
    create_bullet_slide(15, "WORKFLOW STEP 5: MULTICLASS BRAND ATTRIBUTION", s15_data)

    s16_data = [
        {"bold_lead": "Comprehensive Evaluation Metrics:", "text": ""},
        {"bold_lead": "Binary Metrics:", "text": ""},
        {"bold_lead": "F1 Score & MCC:", "text": "Matthews Correlation Coefficient (MCC) accounts for heavy class imbalance across all 4 confusion matrix quadrants.", "level": 1},
        {"bold_lead": "ROC AUC:", "text": "Measures class separability independent of fixed decision cut-off thresholds.", "level": 1},
        {"bold_lead": "Multiclass Metrics:", "text": ""},
        {"bold_lead": "Macro F1 vs. Micro F1:", "text": "Macro F1 treats all brands equally (protecting niche targets); Micro F1 measures global volume accuracy.", "level": 1},
        {"bold_lead": "Identification Rate (Id / Rep_TP):", "text": "Measures recall of correct brand attribution strictly on correctly detected phishing pages."}
    ]
    create_bullet_slide(16, "EVALUATING THE OUTPUT: METRICS & BENCHMARKS", s16_data)

    # =========================================================================
    # SLIDE 17: WORKFLOW DIAGRAM SLIDE
    # =========================================================================
    slide17 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide17, "SYSTEM WORKFLOW & ARCHITECTURE")
    
    blocks_data = [
        ("Inputs", "Product/URL\nScreenshots\n(CERT / PP / VP)", RGBColor(140, 140, 140)),
        ("Data Standardization", "Stratified Splitting\n(60 / 20 / 20)\nAugmentation", ORANGE_RUST),
        ("API Gateway", "FastAPI Orchestrator\nBase64 Validation\nPydantic Schema", DARK_BROWN),
        ("Method Containers", "Isolated Models\n(VisualPhishNet,\nPhishpedia, pHash)", ORANGE_RUST),
        ("Output & Persistence", "Binary Threat Verdict\nTarget Attribution\nSQLite / DB Logging", RGBColor(160, 40, 0))
    ]
    
    start_x = 0.8
    block_w = 2.1
    block_h = 4.2
    gap = 0.3
    
    for i, (title_b, desc_b, color_b) in enumerate(blocks_data):
        bx = Inches(start_x + i * (block_w + gap))
        by = Inches(2.0)
        
        shape = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(block_w), Inches(block_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_b
        shape.line.fill.background()
        
        tf_b = shape.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = Inches(0.15)
        
        p_t = tf_b.paragraphs[0]
        p_t.alignment = PP_ALIGN.CENTER
        r_t = p_t.add_run()
        r_t.text = title_b + "\n\n"
        r_t.font.size = Pt(16)
        r_t.font.bold = True
        r_t.font.color.rgb = WHITE
        
        p_d = tf_b.add_paragraph()
        p_d.alignment = PP_ALIGN.CENTER
        r_d = p_d.add_run()
        r_d.text = desc_b
        r_d.font.size = Pt(13)
        r_d.font.color.rgb = WHITE
        
        if i < len(blocks_data) - 1:
            arrow = slide17.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, bx + Inches(block_w + 0.05), Inches(3.8), Inches(0.2), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED_GRAY
            arrow.line.fill.background()
            
    add_footer(slide17, 17)

    # =========================================================================
    # SLIDE 18: COMPARATIVE RESULTS TABLE
    # =========================================================================
    slide18 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide18, "EXPERIMENTAL RESULTS & COMPARATIVE TABLE")
    
    rows, cols = 10, 6
    table_shape = slide18.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.0))
    table = table_shape.table
    
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(1.7)
    table.columns[3].width = Inches(1.7)
    table.columns[4].width = Inches(1.7)
    table.columns[5].width = Inches(1.933)
    
    headers = ["Dataset", "Method", "F1 Micro", "F1 Macro", "MCC", "Identification Rate"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = WHITE
        
    table_data = [
        ("CERT Polska", "VisualPhishNet", "0.3654", "0.1481", "0.0733", "0.7093"),
        ("CERT Polska", "Phishpedia", "0.5482", "0.2621", "0.2013", "0.9845"),
        ("CERT Polska", "Baseline (pHash)", "0.5823", "0.1670", "0.3668", "0.2554"),
        ("VisualPhishNet (VP)", "VisualPhishNet", "0.3924", "0.0047", "0.0694", "0.0037"),
        ("VisualPhishNet (VP)", "Phishpedia", "0.3782", "0.3073", "0.2569", "0.9270"),
        ("VisualPhishNet (VP)", "Baseline (pHash)", "0.7009", "0.4111", "0.5089", "0.5679"),
        ("Phishpedia (PP)", "VisualPhishNet", "0.1003", "0.0334", "0.0111", "1.0000"),
        ("Phishpedia (PP)", "Phishpedia", "0.7691", "0.2894", "0.7384", "0.9154"),
        ("Phishpedia (PP)", "Baseline (pHash)", "0.5467", "0.3591", "0.4939", "0.5687"),
    ]
    
    for i, row_vals in enumerate(table_data):
        for j, val in enumerate(row_vals):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.size = Pt(12)
            r.font.color.rgb = TEXT_BLACK
            if val in ["0.5823", "0.9845", "0.7009", "0.9270", "0.7691", "0.9154"]:
                r.font.bold = True
                r.font.color.rgb = ORANGE_RUST
                
    add_footer(slide18, 18)

    s19_data = [
        {"bold_lead": "Operational Deployment Scenarios:", "text": ""},
        {"bold_lead": "National CSIRT / CERT Triage Pipeline:", "text": "Processing 15,000+ suspicious regional URLs daily with automated priority scoring.", "level": 1},
        {"bold_lead": "Brand Protection & Anti-Phishing Monitoring:", "text": "Enables automated takedown notifications to impersonated institutions.", "level": 1},
        {"bold_lead": "Browser Security Extensions:", "text": "Real-time client-side pHash verification before user credential submission.", "level": 1},
        {"bold_lead": "CERT Polska Real-World Insights:", "text": "Unaugmented, analyst-verified threat data confirmed that baseline perceptual hashing provides superior stability over complex CNNs in real-time triage."}
    ]
    create_bullet_slide(19, "REAL-WORLD APPLICATIONS & CERT CASE STUDY", s19_data)

    s20_data = [
        {"bold_lead": "Shifting the Paradigm of Phishing Defense:", "text": ""},
        {"bold_lead": "Computational Efficiency:", "text": "VisualPhishNet requires 8-10 hours training per dataset; Baseline pHash & Phishpedia complete training in ~20 minutes.", "level": 1},
        {"bold_lead": "Democratization:", "text": "Provides open-source containerized evaluation tools accessible to small SOC teams without heavy GPU clusters.", "level": 1},
        {"bold_lead": "Operational Agility:", "text": "Allows security teams to instantly update reference logo databases without full model retraining.", "level": 1},
        {"bold_lead": "Analyst Fatigue Reduction:", "text": "Automates initial filtering, reserving human investigation for high-confidence brand matches."}
    ]
    create_bullet_slide(20, "INDUSTRY IMPACT & OPERATIONAL EFFICIENCY", s20_data)

    s21_data = [
        {"bold_lead": "What's Next for Visual Phishing Detection?", "text": ""},
        {"bold_lead": "K-Fold Cross-Validation Framework:", "text": "Incorporating rigorous cross-validation to provide formal uncertainty quantification.", "level": 1},
        {"bold_lead": "Multimodal Deep Fusion:", "text": "Integrating visual pHash/Phishpedia features with DOM Graph Neural Networks (GNN) and LLM reasoning.", "level": 1},
        {"bold_lead": "Temporal Campaign Tracking:", "text": "Grouping detection results by time windows to model long-term phishing infrastructure evolution.", "level": 1},
        {"bold_lead": "Adversarial Robustness:", "text": "Enhancing resilience against adversarial visual perturbations (watermarking, subtle color shifts).", "level": 1}
    ]
    create_bullet_slide(21, "FUTURE SCOPE AND POTENTIAL IMPROVEMENTS", s21_data)

    s22_data = [
        {"bold_lead": "Summary of Key Findings:", "text": ""},
        {"bold_lead": "Baseline (pHash) Superiority:", "text": "Perceptual hashing delivers the most robust binary classification across all benchmark datasets.", "level": 1},
        {"bold_lead": "Phishpedia Target Recognition:", "text": "Consistently achieves >0.90 Identification Rate, proving ideal for brand attribution.", "level": 1},
        {"bold_lead": "VisualPhishNet Vulnerability:", "text": "Suffers from feature collapse when trained on synthetic/augmented data distributions.", "level": 1},
        {"bold_lead": "The Recommended Operational Hybrid Architecture:", "text": ""},
        {"bold_lead": "Stage 1 (Rapid Filtering):", "text": "Use Perceptual Hashing (pHash + FAISS) for instant binary phishing detection.", "level": 1},
        {"bold_lead": "Stage 2 (Brand Identification):", "text": "Pass detected threats to Phishpedia for precise impersonation target recognition.", "level": 1}
    ]
    create_bullet_slide(22, "CONCLUSION", s22_data)

    s23_data = [
        {"bold_lead": "Jarczewski, M., Białczak, P., & Mazurczyk, W.", "text": "(2026). Phishing Website Impersonation: Comparative Analysis of Detection and Target Recognition Methods. MDPI Applied Sciences, 16(2), 640. DOI: 10.3390/app16020640"},
        {"bold_lead": "Lin, Y., Liu, R., Divakaran, D. M., et al.", "text": "(2021). Phishpedia: A Hybrid Deep Learning Based Approach to Visually Identify Phishing Webpages. In 30th USENIX Security Symposium (pp. 3793-3810)."},
        {"bold_lead": "Abdelnabi, S., Krombholz, K., & Fritz, M.", "text": "(2020). VisualPhishNet: Zero-Day Phishing Website Detection by Visual Similarity. In ACM Conference on Computer and Communications Security (CCS) (pp. 1633-1650)."},
        {"bold_lead": "Douze, M., Guzhva, A., Deng, C., et al.", "text": "(2025). The Faiss Library. arXiv preprint arXiv:2401.08281."},
        {"bold_lead": "Framework Code Repository:", "text": "https://github.com/Percival33/phish-target-recognition"}
    ]
    create_bullet_slide(23, "REFERENCES", s23_data)

    # SLIDE 24: THANK YOU
    slide24 = prs.slides.add_slide(blank_slide_layout)
    
    ty_box = slide24.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.733), Inches(2.0))
    tf24 = ty_box.text_frame
    p24 = tf24.paragraphs[0]
    p24.alignment = PP_ALIGN.CENTER
    r24 = p24.add_run()
    r24.text = "THANK YOU"
    r24.font.size = Pt(44)
    r24.font.bold = True
    r24.font.color.rgb = TEXT_BLACK
    r24.font.name = "Calibri"
    
    p24_sub = tf24.add_paragraph()
    p24_sub.alignment = PP_ALIGN.CENTER
    r24_sub = p24_sub.add_run()
    r24_sub.text = "\nQuestions & Comments Welcome"
    r24_sub.font.size = Pt(20)
    r24_sub.font.color.rgb = MUTED_GRAY
    r24_sub.font.name = "Calibri"

    add_footer(slide24, 24)

    output_path = os.path.join(os.getcwd(), "Phishing_Website_Impersonation_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
